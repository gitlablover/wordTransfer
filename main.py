# This is a sample Python script.

# Press Shift+F10 to execute it or replace it with your code.
# Press Double Shift to search everywhere for classes, files, tool windows, actions, and settings.
from datetime import datetime

from docx import Document
from pptx.util import Pt, Cm
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_DATA_LABEL_POSITION
import csv
import numpy as np
import numpy as np
import matplotlib

matplotlib.use('TkAgg')
import matplotlib.pyplot as plt

time = []  # line[0]
timeToshow = []
conc = []  # line[1]
temp = []  # line[2]
nd = []  # line[4]
timeToShowLabel = []

def readCsv():
    csv_reader = csv.reader(open("test.csv"))
    # 取每5分钟的数据
    # 从一个非常早的数据开始
    timeFrom = datetime.strptime('01.01.1966 01:01:00', "%d.%m.%Y %H:%M:%S")
    for line in csv_reader:
        if not str(line[0]).startswith('#'):
            conc.append(float(line[1]))
            temp.append(float(line[2]))
            nd.append(float(line[4]))
            timeStr = str(line[0])

            timeStruct = datetime.strptime(timeStr, "%d.%m.%Y %H:%M:%S")
            time.append(timeStruct)
            if (timeStruct - timeFrom).total_seconds() / 60 >= 5:
                timeFrom = timeStruct
                timeToshow.append(timeStruct)
                timeToShowLabel.append(timeStruct.strftime('%H:%M'))




def drawCONC():

    matplotlib.use('TkAgg')


    plt.figure(figsize=(12,6),dpi=100)
    plt.ylabel('CONC')
    plt.grid(axis='both')
    plt.xticks(timeToshow, timeToShowLabel)
    x = time
    y = conc
    plt.plot(x, y)
    plt.savefig('conc.png')
    plt.show()
def drawTemp():

    matplotlib.use('TkAgg')
    plt.figure(figsize=(12, 6), dpi=100)
    plt.grid(axis='both')
    plt.ylabel('T')
    plt.xticks(timeToshow, timeToShowLabel)
    x = time
    y = temp
    plt.plot(x, y)
    plt.savefig('temp.png')
    plt.show()

def drawnd():

    matplotlib.use('TkAgg')
    plt.figure(figsize=(12, 6), dpi=100)
    plt.grid(axis='both')
    plt.ylabel('nD')
    plt.xticks(timeToshow, timeToShowLabel)
    x = time
    y = nd
    plt.plot(x, y)

    plt.savefig('nd.png')
    plt.show()


def createWord():
    """新建文档"""
    document = Document()

    document.add_picture('conc.png', width=Cm(14.0))
    document.add_picture('temp.png', width=Cm(14.0))
    document.add_picture('nd.png', width=Cm(14.0))

    document.save('test.docx')


# Press the green button in the gutter to run the script.
if __name__ == '__main__':
    readCsv()

    drawCONC()
    drawTemp()
    drawnd()
    createWord()

# See PyCharm help at https://www.jetbrains.com/help/pycharm/
